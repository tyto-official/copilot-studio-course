#!/usr/bin/env python3
"""Build 'Begrepp och Teori' presentation for Copilot Studio course.

Usage:
    python build_presentation.py
    python build_presentation.py --pdf
"""

import argparse
import os
import shutil
import subprocess
import sys
from datetime import datetime
from pathlib import Path
import time

# Prefer the project's bundled packages so the script still runs even if the
# local virtualenv launcher is broken.
ROOT = Path(__file__).resolve().parent
LOCAL_SITE_PACKAGES = ROOT / "venv" / "Lib" / "site-packages"
if LOCAL_SITE_PACKAGES.exists():
    sys.path.insert(0, str(LOCAL_SITE_PACKAGES))

from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────
BILDER = ROOT / "bilder"
VIDEOS = ROOT / "videos"
AUDIO = ROOT / "ljud"
CHAP = ROOT.parent / "copilot-studio-course" / "docs" / "assets" / "images" / "chap"
OUTPUT = ROOT / "Begrepp-och-teori.pptx"
OUTPUT_PDF = ROOT / "Begrepp-och-teori.pdf"

# ── Image catalog ──────────────────────────────────────────────────────────
IMG = {
    # bilder/
    "hit100m":       BILDER / "hit-100m-users.jpg",
    "fore":          BILDER / "fore.jpg",
    "efter":         BILDER / "efter.png",
    "joel":          BILDER / "joel.png",
    "sverige":       BILDER / "sverige.jpg",
    "joel_sverige":  BILDER / "joel-med-sverige.png",
    "nyheter":       BILDER / "nyheter.png",
    "bocker":        BILDER / "böcker.jpg",
    "pdf":           BILDER / "pdf.png",
    "youtube":       BILDER / "youtube-transkibering.png",
    "old_chatgpt":   BILDER / "old-chatgpt.png",
    "intro_chatgpt": BILDER / "intoducion-chatgpt.png",
    "nvidia":        BILDER / "nvidia-stock.png",
    "gpu_price":     BILDER / "gpu-price.png",
    "ranking":       BILDER / "ranking.png",
    "token":         BILDER / "token.png",
    "token_id":      BILDER / "token-id.png",
    "workflows":     BILDER / "workflows.jpeg",
    "ai_agent":      BILDER / "ai-agnet.jpeg",
    "forsta_bakrund": BILDER / "forsta_bakrund.jpg",
    "forsta_byggblocken": BILDER / "forsta_byggblocken.jpg",
    "forsta_anvanding": BILDER / "forsta_anvanding.jpg",
    # chap/
    "bil1":          CHAP / "bil1.jpg",
    "bil2":          CHAP / "bil2.jpg",
    "bil3":          CHAP / "bil3.jpg",
    "bil4":          CHAP / "bil4.jpg",
    "bil5":          CHAP / "bil5.png",
    "chunks":        CHAP / "2.png",
    "vektor":        CHAP / "3.jpeg",
    "hitta":         CHAP / "4.jpeg",
    "kontext":       CHAP / "7.jpeg",
    "resoning":      CHAP / "resoning.png",
    "valjmodell":    CHAP / "valjmodell.png",
    "verktyg":       CHAP / "verktyg.png",
    "ollama":        CHAP / "ollama-ministral-3b.png",
}

MEDIA = {
    "video_1": VIDEOS / "PgDwfd1p.mp4",
    "video_2": VIDEOS / "HpJVHG4G.mp4",
    "video_3": VIDEOS / "ngTyKJHd.mp4",
    "audio_1": AUDIO / "elevenlabs.mp3",
}

# ── Design tokens ──────────────────────────────────────────────────────────
FT = "Bahnschrift SemiBold"
FB = "Aptos"


class C:
    """Color palette."""
    NAVY      = RGBColor(20, 33, 69)
    CANVAS    = RGBColor(246, 248, 253)
    WHITE     = RGBColor(255, 255, 255)
    PANEL_ALT = RGBColor(232, 238, 250)
    TEXT      = RGBColor(29, 41, 79)
    MUTED     = RGBColor(95, 110, 146)
    ACCENT    = RGBColor(40, 64, 127)
    ACCENT2   = RGBColor(50, 77, 147)
    BORDER    = RGBColor(204, 215, 239)
    DARK      = RGBColor(35, 55, 108)
    DARK_LINE = RGBColor(88, 118, 196)
    MID_BLUE  = RGBColor(98, 128, 189)
    LIGHT_BLU = RGBColor(101, 142, 206)
    DEEP_BLUE = RGBColor(75, 98, 158)
    WARM_BG   = RGBColor(254, 243, 199)
    WARM_LINE = RGBColor(245, 158, 11)
    STEP_LBL  = RGBColor(154, 52, 18)


# ── Helpers ────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill, line=None, lw=0.75, rounded=True, radius=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(st, Pt(l), Pt(t), Pt(w), Pt(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.fill.solid()
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    if radius is not None:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape


def txt(slide, l, t, w, h, text, font=FB, size=18, color=C.TEXT,
        bold=False, align=PP_ALIGN.LEFT, margin=0):
    box = slide.shapes.add_textbox(Pt(l), Pt(t), Pt(w), Pt(h))
    tf = box.text_frame
    tf.margin_left = Pt(margin)
    tf.margin_right = Pt(margin)
    tf.margin_top = Pt(margin)
    tf.margin_bottom = Pt(margin)
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        p.alignment = align
    return box


def pic(slide, key, l, t, w, h):
    path = IMG.get(key)
    if not path or not path.exists():
        rect(slide, l, t, w, h, C.PANEL_ALT, C.BORDER)
        txt(slide, l, t + h // 2 - 10, w, 20, f"[{key}]",
            size=10, color=C.MUTED, align=PP_ALIGN.CENTER)
        return None
    im = Image.open(path)
    iw, ih = im.size
    im.close()
    scale = min(w / iw, h / ih)
    tw, th = iw * scale, ih * scale
    tl, tt = l + (w - tw) / 2, t + (h - th) / 2
    return slide.shapes.add_picture(str(path), Pt(tl), Pt(tt), Pt(tw), Pt(th))


def chrome(slide, section, title, accent=C.ACCENT):
    rect(slide, 0, 0, 960, 540, C.CANVAS, rounded=False)
    rect(slide, 0, 0, 960, 12, accent, rounded=False)
    txt(slide, 60, 28, 280, 24, section, FB, 12, accent, bold=True)
    txt(slide, 60, 52, 820, 42, title, FT, 27, C.TEXT, bold=True)
    rect(slide, 60, 102, 840, 1.5, C.BORDER, rounded=False)
    txt(slide, 60, 510, 180, 14, "Copilot Studio", FB, 9, C.MUTED)


def fnum(slide, n):
    rect(slide, 885, 500, 34, 24, C.WHITE, C.BORDER, radius=0.15)
    txt(slide, 885, 503, 34, 18, str(n), FB, 10, C.TEXT,
        bold=True, align=PP_ALIGN.CENTER)


def step(slide, l, t, w, h, num, title, body):
    rect(slide, l, t, w, h, C.WHITE, C.BORDER)
    rect(slide, l + 16, t + 14, 32, 24, C.ACCENT, radius=0.15)
    txt(slide, l + 16, t + 17, 32, 18, num, FB, 11, C.NAVY,
        bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l + 58, t + 10, w - 72, 24, title, FT, 15, C.TEXT, bold=True)
    b = txt(slide, l + 18, t + 40, w - 36, h - 54, body, FB, 11.5, C.MUTED)
    for p in b.text_frame.paragraphs:
        p.space_after = Pt(4)


def step_illustrated(slide, l, t, w, h, num, title, body, image_key):
    rect(slide, l, t, w, h, C.WHITE, C.BORDER)
    rect(slide, l + 16, t + 14, 32, 24, C.ACCENT, radius=0.15)
    txt(slide, l + 16, t + 17, 32, 18, num, FB, 11, C.NAVY,
        bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l + 58, t + 10, w - 72, 24, title, FT, 15, C.TEXT, bold=True)
    b = txt(slide, l + 18, t + 40, w - 36, 106, body, FB, 11.5, C.MUTED)
    for p in b.text_frame.paragraphs:
        p.space_after = Pt(4)
    pic(slide, image_key, l + 18, t + h - 106, w - 36, 88)


def bullets(slide, l, t, w, h, lines, size=18, color=None):
    if color is None:
        color = C.TEXT
    box = slide.shapes.add_textbox(Pt(l), Pt(t), Pt(w), Pt(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u2022 {line}"
        run.font.name = FB
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def pick_output_path(preferred):
    """Use the preferred output unless it is locked, then fall back."""
    if not preferred.exists():
        return preferred
    try:
        with preferred.open("ab"):
            pass
        return preferred
    except PermissionError:
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        fallback = preferred.with_name(f"{preferred.stem}-{stamp}{preferred.suffix}")
        print(f"Output locked, using: {fallback}")
        return fallback


def export_pdf(pptx_path, pdf_path):
    """Export the generated presentation to PDF via local PowerPoint."""
    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:
        print(f"PDF export skipped: missing dependency ({exc}).")
        return False

    app = None
    presentation = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = 1
        presentation = app.Presentations.Open(
            str(pptx_path),
            ReadOnly=1,
            Untitled=0,
            WithWindow=0,
        )
        presentation.SaveAs(str(pdf_path), 32)
        print(f"Created: {pdf_path}")
        return True
    except Exception as exc:
        print(f"PDF export skipped: {exc}")
        return False
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def find_com_python():
    """Return a Python executable likely to have pywin32 installed."""
    candidates = [
        Path(r"C:\Users\zille\anaconda3\python.exe"),
    ]
    for env_name in ("CONDA_PYTHON_EXE",):
        value = os.environ.get(env_name)
        if value:
            candidates.append(Path(value))

    which_python = shutil.which("python")
    if which_python:
        candidates.append(Path(which_python))

    seen = set()
    for candidate in candidates:
        candidate = candidate.resolve()
        if candidate in seen:
            continue
        seen.add(candidate)
        if candidate.exists() and candidate != Path(sys.executable).resolve():
            return candidate
    return None


def populate_media_slides_via_external_python(pptx_path, exc):
    """Retry media embedding using a Python installation with pywin32."""
    com_python = find_com_python()
    if com_python is None:
        print(f"Media embedding skipped: missing dependency ({exc}).")
        return False

    cmd = [str(com_python), str(Path(__file__).resolve()), "--media-pass", str(pptx_path)]
    result = subprocess.run(cmd, cwd=str(ROOT))
    if result.returncode == 0:
        print(f"Media embedding completed via: {com_python}")
        return True

    print(f"Media embedding skipped: helper Python failed ({com_python}).")
    return False


def populate_media_slides(pptx_path, allow_external_python=True):
    """Populate the media slides using PowerPoint so video/audio embed cleanly."""
    try:
        import pythoncom
        import win32com.client
        from win32com.client import constants as c
    except ImportError as exc:
        if allow_external_python:
            return populate_media_slides_via_external_python(pptx_path, exc)
        print(f"Media embedding skipped: missing dependency ({exc}).")
        return False

    app = None
    presentation = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = 1
        last_error = None
        for _ in range(5):
            try:
                presentation = app.Presentations.Open(
                    str(pptx_path),
                    ReadOnly=0,
                    Untitled=0,
                    WithWindow=0,
                )
                break
            except Exception as exc:
                last_error = exc
                time.sleep(0.8)
        if presentation is None:
            raise last_error

        video_slide = presentation.Slides(6)
        video_shapes = []
        video_area = (60, 72, 840, 420)
        for key in ("video_1", "video_2", "video_3"):
            movie_file = MEDIA[key]
            if not Path(movie_file).exists():
                continue
            media_shape = video_slide.Shapes.AddMediaObject2(
                str(movie_file), False, True, *video_area
            )
            media_shape.MediaFormat.Muted = True
            media_shape.MediaFormat.Volume = 0
            media_shape.AnimationSettings.PlaySettings.PlayOnEntry = True
            media_shape.AnimationSettings.PlaySettings.RewindMovie = True
            media_shape.AnimationSettings.PlaySettings.HideWhileNotPlaying = True
            video_shapes.append(media_shape)

        if video_shapes:
            first_video = video_shapes[0]
            first_video.AnimationSettings.PlaySettings.HideWhileNotPlaying = False

        for idx, media_shape in enumerate(video_shapes[:-1]):
            fade_out = video_slide.TimeLine.MainSequence.AddEffect(
                media_shape, c.msoAnimEffectFade
            )
            fade_out.Exit = True
            fade_out.Timing.Duration = 0.35

            next_shape = video_shapes[idx + 1]
            if idx + 1 < len(video_shapes) - 1:
                fade_in = video_slide.TimeLine.MainSequence.AddEffect(
                    next_shape, c.msoAnimEffectFade
                )
                fade_in.Timing.Duration = 0.35

        if video_shapes:
            seq = video_slide.TimeLine.MainSequence
            play_effects = {}
            entry_effects = {}
            exit_effects = {}
            for i in range(1, seq.Count + 1):
                effect = seq.Item(i)
                shape_name = effect.Shape.Name
                if effect.EffectType == 83:
                    play_effects[shape_name] = effect
                elif effect.EffectType == 10 and effect.Exit:
                    exit_effects[shape_name] = effect
                elif effect.EffectType == 10:
                    entry_effects[shape_name] = effect

            last_name = video_shapes[-1].Name
            if last_name in entry_effects:
                entry_effects[last_name].Delete()
                del entry_effects[last_name]

            ordered_effects = []
            first_name = video_shapes[0].Name
            if first_name in play_effects:
                ordered_effects.append(play_effects[first_name])

            for idx, media_shape in enumerate(video_shapes[:-1]):
                shape_name = media_shape.Name
                next_name = video_shapes[idx + 1].Name
                if shape_name in exit_effects:
                    ordered_effects.append(exit_effects[shape_name])
                if next_name in entry_effects:
                    ordered_effects.append(entry_effects[next_name])
                if next_name in play_effects:
                    ordered_effects.append(play_effects[next_name])

            for index, effect in enumerate(ordered_effects, start=1):
                effect.MoveTo(index)

            if ordered_effects:
                ordered_effects[0].Timing.TriggerType = c.msoAnimTriggerAfterPrevious
                ordered_effects[0].Timing.TriggerDelayTime = 0

            for effect in ordered_effects[1:]:
                effect.Timing.TriggerType = c.msoAnimTriggerAfterPrevious
                effect.Timing.TriggerDelayTime = 0
                if effect.EffectType == 10:
                    effect.Timing.Duration = 0.35

        audio_slide = presentation.Slides(7)
        if MEDIA["audio_1"].exists():
            audio_shape = audio_slide.Shapes.AddMediaObject2(
                str(MEDIA["audio_1"]), False, True, 428, 225, 104, 104
            )
            audio_shape.MediaFormat.Volume = 1

        presentation.Save()
        return True
    except Exception as exc:
        if allow_external_python:
            return populate_media_slides_via_external_python(pptx_path, exc)
        print(f"Media embedding skipped: {exc}")
        return False
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


# ── Build ──────────────────────────────────────────────────────────────────

def build(make_pdf=False):
    prs = Presentation()
    prs.slide_width = Pt(960)
    prs.slide_height = Pt(540)
    n = 0
    pptx_output = pick_output_path(OUTPUT)
    pdf_output = pick_output_path(OUTPUT_PDF)
    if pptx_output.stem != OUTPUT.stem:
        pdf_output = pdf_output.with_name(f"{pptx_output.stem}.pdf")

    # ── 1. Title ───────────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    rect(s, 0, 0, 960, 540, C.NAVY, rounded=False)
    rect(s, 0, 0, 960, 14, C.ACCENT2, rounded=False)
    txt(s, 60, 58, 260, 24, "Del 1", FB, 16, C.PANEL_ALT, bold=True)
    txt(s, 60, 88, 420, 74, "Begrepp och teori", FT, 34, C.WHITE, bold=True)
    txt(s, 60, 178, 380, 86,
        "Den teoretiska grunden bakom\nspr\u00e5kmodeller, kontext, RAG\noch AI-agenter innan vi b\u00f6rjar bygga.",
        FB, 20, C.WHITE)
    txt(s, 60, 458, 260, 20, "Microsoft Copilot Studio",
        FB, 13, C.PANEL_ALT, bold=True)
    rect(s, 510, 52, 366, 194, C.DARK, C.DARK_LINE)
    rect(s, 510, 266, 366, 194, C.DARK, C.DARK_LINE)
    pic(s, "hit100m", 528, 72, 154, 154)
    pic(s, "token_id", 704, 72, 154, 154)
    pic(s, "pdf", 528, 286, 154, 154)
    pic(s, "ranking", 704, 286, 154, 154)

    # ── 2. Varf\u00f6r? ───────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Varf\u00f6r denna teoridel", "Tre anledningar till att vi b\u00f6rjar h\u00e4r",
           C.ACCENT2)
    step_illustrated(s, 60, 136, 250, 286, "1", "F\u00f6rst\u00e5 bakgrunden",
                     "Du f\u00e5r en tydligare bild av vad som faktiskt sker bakom kulisserna i moderna AI-system.",
                     "forsta_bakrund")
    step_illustrated(s, 330, 136, 250, 286, "2", "F\u00f6rst\u00e5 byggblocken",
                     "Du ser vilka komponenter som beh\u00f6vs f\u00f6r att bygga effektiva agenter i Copilot Studio.",
                     "forsta_byggblocken")
    step_illustrated(s, 600, 136, 250, 286, "3", "F\u00f6rst\u00e5 avv\u00e4gningarna",
                     "Du f\u00e5r b\u00e4ttre intuition f\u00f6r m\u00f6jligheter, begr\u00e4nsningar och designval n\u00e4r vi senare bygger.",
                     "forsta_anvanding")
    fnum(s, n)

    # ── 3. Spr\u00e5kmodeller intro ───────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Spr\u00e5kmodeller",
           "Den viktigaste AI-teknologin i den h\u00e4r kursen", C.ACCENT)
    txt(s, 60, 138, 300, 150,
        "Spr\u00e5kmodellen \u00e4r ofta hj\u00e4rnan i agenten. Det var ocks\u00e5 spr\u00e5kmodeller som gjorde AI till ett massfenomen n\u00e4r ChatGPT lanserades i slutet av november 2022.",
        FB, 22, C.TEXT)
    bullets(s, 60, 312, 300, 96, [
        "AI-hajpen tog fart via chattgr\u00e4nssnittet.",
        "Det drog i sin tur upp investeringar, forskning och f\u00f6rv\u00e4ntningar.",
    ], 16)
    rect(s, 394, 136, 506, 332, C.WHITE, C.BORDER)
    pic(s, "hit100m", 414, 156, 466, 292)
    fnum(s, n)

    # ── 4. Bildgenerering: f\u00f6re / efter ───────────────────────────────
    n += 1
    s = new_slide(prs)
    rect(s, 0, 0, 960, 540, C.NAVY, rounded=False)
    txt(s, 60, 28, 300, 24, "Bildgenerering", FB, 13, C.PANEL_ALT, bold=True)
    txt(s, 60, 54, 520, 34, "Originalbild och resultat", FT, 26, C.WHITE, bold=True)
    txt(s, 80, 94, 360, 20, "Originalbild", FB, 12, C.PANEL_ALT, bold=True)
    txt(s, 520, 94, 360, 20, "Efter prompt: addera en tiger", FB, 12, C.PANEL_ALT, bold=True)
    pic(s, "fore", 40, 120, 420, 360)
    pic(s, "efter", 500, 120, 420, 360)
    fnum(s, n)

    # ── 5. Kombinera signaler ─────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Kombinera signaler",
           "En modell kan bygga nytt inneh\u00e5ll fr\u00e5n flera visuella referenser",
           C.ACCENT)
    for key, label, x in [
        ("joel", "Bild p\u00e5 mig", 60),
        ("sverige", "Sverigetr\u00f6ja", 340),
        ("joel_sverige", "Genererat resultat", 620),
    ]:
        rect(s, x, 142, 240, 250, C.WHITE, C.BORDER)
        pic(s, key, x + 16, 158, 208, 208)
        txt(s, x + 16, 372, 208, 18, label, FB, 11, C.MUTED,
            align=PP_ALIGN.CENTER)
    fnum(s, n)

    # ── 6. Videoexempel ───────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    rect(s, 0, 0, 960, 540, C.NAVY, rounded=False)
    rect(s, 0, 0, 960, 12, C.ACCENT2, rounded=False)
    txt(s, 60, 26, 240, 20, "Videoexempel", FB, 12, C.PANEL_ALT, bold=True)
    fnum(s, n)

    # ── 7. Ljudexempel ────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Ljud", "Kort ljudsnutt", C.ACCENT)
    txt(s, 60, 146, 840, 26,
        "Exempel med ElevenLabs-ljud. Klicka p\u00e5 ljudikonen f\u00f6r att spela upp klippet.",
        FB, 16, C.TEXT, align=PP_ALIGN.CENTER)
    txt(s, 60, 346, 840, 22, "ElevenLabs", FT, 20, C.ACCENT2, bold=True,
        align=PP_ALIGN.CENTER)
    fnum(s, n)

    # ── 8. Sektionsstart ──────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    rect(s, 0, 0, 960, 540, C.NAVY, rounded=False)
    rect(s, 0, 0, 960, 14, C.ACCENT2, rounded=False)
    txt(s, 60, 160, 840, 56, "Hur bygger man en språkmodell?", FT, 34, C.WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, 120, 236, 720, 52,
        "Nu går vi från exempel till själva byggprocessen: data in, träning och därefter körning.",
        FB, 20, C.WHITE, align=PP_ALIGN.CENTER)
    fnum(s, n)

    # ── 9. Datainsamling ──────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Del 1", "Datainsamling", C.ACCENT2)
    txt(s, 60, 132, 320, 88,
        "F\u00f6r att bygga en spr\u00e5kmodell kr\u00e4vs extrema m\u00e4ngder data, framf\u00f6r allt i textformat. Det kan vara nyheter, b\u00f6cker, PDF:er och transkriberat ljud eller video.",
        FB, 20, C.TEXT)
    for key, label, x, y in [
        ("nyheter", "Nyheter", 400, 130),
        ("bocker", "B\u00f6cker", 650, 130),
        ("pdf", "PDF-dokument", 400, 302),
        ("youtube", "Transkribering / video", 650, 302),
    ]:
        rect(s, x, y, 230, 150, C.WHITE, C.BORDER)
        pic(s, key, x + 14, y + 14, 202, 122)
        txt(s, x + 14, y + 130, 202, 16, label, FB, 10.5, C.MUTED,
            align=PP_ALIGN.CENTER)
    fnum(s, n)

    # ── 8. Grundtr\u00e4ning ─────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Del 1", "Grundtr\u00e4ning", C.ACCENT)
    step(s, 60, 152, 180, 190, "1", "Se text",
         "Modellen f\u00e5r se enorma m\u00e4ngder tr\u00e4ningsdata.")
    step(s, 260, 152, 180, 190, "2", "Gissa n\u00e4sta",
         "Den f\u00f6rs\u00f6ker f\u00f6ruts\u00e4ga vad som mest sannolikt kommer h\u00e4rn\u00e4st.")
    step(s, 460, 152, 180, 190, "3", "M\u00e4t fel",
         "Systemet j\u00e4mf\u00f6r gissningen med r\u00e4tt forts\u00e4ttning.")
    step(s, 660, 152, 180, 190, "4", "Justera vikter",
         "Parametrarna finjusteras lite grand varje varv.")
    rect(s, 60, 368, 780, 90, C.PANEL_ALT, C.BORDER, radius=0.12)
    txt(s, 80, 386, 740, 54,
        "Denna process upprepas p\u00e5 enorma m\u00e4ngder text. I b\u00f6rjan gissar modellen ofta fel, men vikterna justeras steg f\u00f6r steg tills modellen gradvis blir b\u00e4ttre.",
        FB, 14, C.MUTED)
    fnum(s, n)

    # ── 9. Bil-analogi ────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Intuition", "Sj\u00e4lvk\u00f6rande bil genom hinderbana", C.ACCENT2)
    for key, label, x in [
        ("bil1", "B\u00f6rjan: kaos", 60),
        ("bil2", "Tidigt: b\u00f6rjar hitta r\u00e4tt", 290),
        ("bil3", "Senare: fler klarar sv\u00e4ngen", 520),
        ("bil4", "Slutet: klarar banan", 750),
    ]:
        rect(s, x, 142, 180, 220, C.WHITE, C.BORDER)
        pic(s, key, x + 8, 150, 164, 164)
        txt(s, x + 8, 320, 164, 30, label, FB, 10, C.MUTED,
            align=PP_ALIGN.CENTER)
    txt(s, 60, 396, 870, 40,
        "Fr\u00e5n m\u00e5nga fel i b\u00f6rjan till en modell som gradvis l\u00e4r sig k\u00f6ra b\u00e4ttre.",
        FB, 16, C.MUTED)
    fnum(s, n)

    # ── 10. Bil: resultat ─────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Intuition",
           "Efter tr\u00e4ningen anv\u00e4nds modellen i nya situationer", C.ACCENT)
    rect(s, 60, 130, 840, 350, C.WHITE, C.BORDER)
    pic(s, "bil5", 80, 148, 800, 314)
    fnum(s, n)

    # ── 11. Vikt-intuition ────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Intuition", "En enkel intuition f\u00f6r vad en vikt g\u00f6r",
           C.ACCENT2)
    # Left panel
    rect(s, 60, 130, 400, 260, C.WHITE, C.BORDER)
    txt(s, 82, 148, 360, 22, "F\u00f6renklad ekvation", FT, 14,
        C.STEP_LBL, bold=True)
    txt(s, 82, 180, 360, 36, "y = w \u00d7 x", FT, 30, C.TEXT, bold=True)
    txt(s, 82, 226, 360, 20, "x = inputsignal (t.ex. texten)", FB, 13, C.MUTED)
    txt(s, 82, 250, 360, 20, "w = vikten (justeras under tr\u00e4ningen)",
        FB, 13, C.MUTED)
    txt(s, 82, 274, 360, 20, "y = output (po\u00e4ng f\u00f6r ett m\u00f6jligt ord)",
        FB, 13, C.MUTED)
    # Right panel
    rect(s, 500, 130, 400, 260, C.WARM_BG, C.WARM_LINE)
    txt(s, 522, 148, 360, 22, 'Exempel: ordet "frukost"', FT, 14,
        C.STEP_LBL, bold=True)
    txt(s, 522, 182, 360, 18, "Om vikten \u00e4r f\u00f6r l\u00e5g:", FB, 13, C.MUTED)
    txt(s, 522, 206, 360, 28, "y = 0.2 \u00d7 10 = 2", FT, 22, C.TEXT,
        bold=True)
    txt(s, 522, 240, 360, 18,
        'L\u00e5g po\u00e4ng \u2192 modellen v\u00e4ljer troligen inte "frukost"',
        FB, 12, C.MUTED)
    txt(s, 522, 274, 360, 18, "Om vikten justeras upp:", FB, 13, C.MUTED)
    txt(s, 522, 298, 360, 28, "y = 0.8 \u00d7 10 = 8", FT, 22, C.TEXT,
        bold=True)
    txt(s, 522, 332, 360, 18,
        'H\u00f6g po\u00e4ng \u2192 "frukost" blir det troliga valet',
        FB, 12, C.MUTED)
    fnum(s, n)

    # ── 12. Fr\u00e5n spr\u00e5kmodell till assistent ────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Del 2", "Fr\u00e5n spr\u00e5kmodell till assistent", C.ACCENT2)
    step(s, 60, 170, 250, 220, "1", "Bred spr\u00e5kf\u00f6rst\u00e5else",
         "Efter grundtr\u00e4ningen kan modellen spr\u00e5k, struktur och mycket generell kunskap.")
    step(s, 330, 170, 250, 220, "2", "Instruktionsf\u00f6ljsamhet",
         "Genom nya tr\u00e4ningsfaser l\u00e4r den sig att f\u00f6lja fr\u00e5gor b\u00e4ttre, h\u00e5lla r\u00e4tt ton och bli mer hj\u00e4lpsam.")
    rect(s, 610, 152, 250, 240, C.WHITE, C.BORDER)
    pic(s, "old_chatgpt", 628, 176, 96, 120)
    pic(s, "intro_chatgpt", 734, 176, 108, 120)
    txt(s, 628, 314, 214, 46,
        "Po\u00e4ngen h\u00e4r \u00e4r inte bara mer kunskap, utan b\u00e4ttre beteende som assistent.",
        FB, 14, C.MUTED)
    fnum(s, n)

    # ── 13. Feedback och finjustering ─────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Del 3", "Feedback och finjustering", C.ACCENT)
    step(s, 60, 170, 250, 220, "1", "J\u00e4mf\u00f6r svar",
         "Olika svar kan st\u00e4llas mot varandra f\u00f6r att avg\u00f6ra vilket som verkar mest hj\u00e4lpsamt och s\u00e4kert.")
    step(s, 330, 170, 250, 220, "2", "Styr beteendet",
         "M\u00e5let \u00e4r att forma b\u00e4ttre beteende, inte att modellen n\u00f6dv\u00e4ndigtvis l\u00e4r sig helt nya fakta.")
    rect(s, 610, 170, 250, 220, C.PANEL_ALT, C.BORDER, radius=0.12)
    txt(s, 628, 192, 214, 22, "Sammanfattning", FT, 14, C.ACCENT2, bold=True)
    txt(s, 628, 224, 214, 140,
        "Del 1 bygger hj\u00e4rnan.\nDel 2 l\u00e4r modellen svara som en assistent.\nDel 3 finslipar beteendet ytterligare.",
        FB, 13, C.TEXT)
    fnum(s, n)

    # ── 14. K\u00f6rning / Inference ──────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "K\u00f6rning", "Vad h\u00e4nder n\u00e4r modellen k\u00f6rs?", C.ACCENT2)
    step(s, 60, 152, 180, 190, "1", "Fr\u00e5ga in",
         'Du skickar in din fr\u00e5ga, t.ex. "Vad \u00e4r en agent?"')
    step(s, 260, 152, 180, 190, "2", "Anv\u00e4nd vikterna",
         "Modellen r\u00e4knar ut vad som mest sannolikt kommer h\u00e4rn\u00e4st.")
    step(s, 460, 152, 180, 190, "3", "Bygg svaret",
         "Svaret byggs upp steg f\u00f6r steg, ord f\u00f6r ord.")
    step(s, 660, 152, 180, 190, "4", "Svar ut",
         "Det f\u00e4rdiga svaret skickas tillbaka till dig.")
    rect(s, 60, 368, 780, 90, C.PANEL_ALT, C.BORDER, radius=0.12)
    txt(s, 80, 386, 740, 54,
        "Det \u00e4r de tr\u00e4nade vikterna som g\u00f6r arbetet. Tr\u00e4ningen st\u00e4ller in vikterna \u2013 k\u00f6rningen anv\u00e4nder dem. Kraftiga GPU:er kr\u00e4vs eftersom ber\u00e4kningarna \u00e4r tunga.",
        FB, 14, C.MUTED)
    fnum(s, n)

    # ── 15. GPU:er ────────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "GPU:er", "K\u00f6rning kr\u00e4ver tung ber\u00e4kningskraft", C.ACCENT)
    rect(s, 60, 138, 500, 192, C.WHITE, C.BORDER)
    pic(s, "nvidia", 80, 154, 460, 160)
    txt(s, 60, 352, 500, 56,
        "Nvidia blev en tydlig symbol f\u00f6r AI-boomen, eftersom moderna spr\u00e5kmodeller ofta k\u00f6rs p\u00e5 GPU:er.",
        FB, 18, C.TEXT)
    rect(s, 590, 138, 300, 270, C.WHITE, C.BORDER)
    pic(s, "gpu_price", 608, 156, 264, 234)
    fnum(s, n)

    # ── 16. Kostnad ───────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Kostnad", "AI-infrastruktur \u00e4r dyr", C.ACCENT2)
    rect(s, 60, 138, 470, 242, C.WHITE, C.BORDER)
    pic(s, "gpu_price", 78, 156, 434, 206)
    txt(s, 560, 150, 300, 176,
        "Kraftiga GPU:er \u00e4r en central kostnadsfaktor. Det p\u00e5verkar varf\u00f6r vissa modeller k\u00f6rs lokalt, andra i molnet och varf\u00f6r stora modeller ofta blir dyrare att anv\u00e4nda i praktiken.",
        FB, 20, C.TEXT)
    fnum(s, n)

    # ── 17. Modellval ─────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Modellval", "Olika typer av spr\u00e5kmodeller", C.ACCENT)
    txt(s, 60, 140, 260, 126,
        "Modeller skiljer sig i kapacitet, kostnad, \u00f6ppenhet och var de k\u00f6rs. Modellval \u00e4r alltid en designfr\u00e5ga, inte bara en kvalitetsfr\u00e5ga.",
        FB, 20, C.TEXT)
    rect(s, 344, 136, 556, 146, C.WHITE, C.BORDER)
    pic(s, "ranking", 362, 148, 520, 122)
    rect(s, 344, 300, 556, 146, C.WHITE, C.BORDER)
    pic(s, "valjmodell", 362, 312, 520, 122)
    fnum(s, n)

    # ── 18. Open / Closed source ──────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "\u00d6ppenhet", "Open source och closed source", C.ACCENT2)
    # Open source
    rect(s, 60, 152, 370, 290, C.WHITE, C.BORDER)
    txt(s, 82, 168, 326, 22, "Open source", FT, 16, C.ACCENT, bold=True)
    txt(s, 82, 198, 326, 54,
        "G\u00e5r ofta att ladda ner och k\u00f6ra sj\u00e4lv. Du kan se hur modellen \u00e4r uppbyggd och arbeta med den lokalt.",
        FB, 14, C.TEXT)
    pic(s, "ollama", 82, 264, 326, 156)
    # Closed source
    rect(s, 470, 152, 370, 290, C.WHITE, C.BORDER)
    txt(s, 492, 168, 326, 22, "Closed source", FT, 16, C.ACCENT, bold=True)
    txt(s, 492, 198, 326, 54,
        "Anv\u00e4nds via en tj\u00e4nst. Du kan anv\u00e4nda modellen, men inte ladda ner den eller se exakt hur vikterna ser ut.",
        FB, 14, C.TEXT)
    pic(s, "ranking", 492, 264, 326, 156)
    fnum(s, n)

    # ── 19. Var bor modellen? ─────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Drift", "Var bor modellen?", C.ACCENT)
    step(s, 60, 160, 260, 240, "1", "Lokal modell",
         "Ligger p\u00e5 din egen dator och k\u00f6rs av din egen h\u00e5rdvara. Fungerar f\u00f6r mindre open source-modeller.")
    step(s, 340, 160, 260, 240, "2", "Molnmodell",
         "Ligger hos leverant\u00f6ren och k\u00f6rs p\u00e5 deras servrar. Copilot Studio anv\u00e4nder Azure OpenAI i Microsofts moln.")
    step(s, 620, 160, 260, 240, "3", "S\u00e4kerhet",
         "Du beh\u00f6ver lita p\u00e5 leverant\u00f6ren. Microsoft k\u00f6r modellen i Azure \u2013 inte via OpenAI:s publika ChatGPT-tj\u00e4nst.")
    fnum(s, n)

    # ── 20. Tokens ────────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    rect(s, 0, 0, 960, 540, C.NAVY, rounded=False)
    txt(s, 60, 28, 300, 24, "Tokens", FB, 13, C.PANEL_ALT, bold=True)
    txt(s, 60, 54, 420, 34, "Text uppdelad i tokens", FT, 26, C.WHITE,
        bold=True)
    pic(s, "token", 50, 102, 860, 390)
    fnum(s, n)

    # ── 21. Token-ID ──────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    rect(s, 0, 0, 960, 540, C.NAVY, rounded=False)
    txt(s, 60, 28, 300, 24, "Tokens", FB, 13, C.PANEL_ALT, bold=True)
    txt(s, 60, 54, 420, 34, "Samma text som token-ID:n", FT, 26, C.WHITE,
        bold=True)
    pic(s, "token_id", 50, 102, 860, 390)
    fnum(s, n)

    # ── 22. Token-representation ──────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Representation", "Samma inneh\u00e5ll p\u00e5 tre niv\u00e5er", C.ACCENT2)
    for label, content, y in [
        ("Vanlig text", "OpenAI!", 130),
        ("Tokens (f\u00f6renklat)", "Open  |  AI  |  !", 215),
        ("Token-ID:n", "[6447, 17527, 0]", 300),
        ("Bitar", "01001111 01110000 01100101 01101110 01000001 01001001 00100001", 385),
    ]:
        rect(s, 60, y, 840, 65, C.WHITE, C.BORDER)
        txt(s, 82, y + 8, 200, 18, label, FB, 12, C.STEP_LBL, bold=True)
        txt(s, 82, y + 30, 796, 24, content, FT, 18, C.TEXT, bold=True)
    fnum(s, n)

    # ── 23. Prompts ───────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Prompts", "System Prompt och User Prompt", C.ACCENT)
    step(s, 60, 162, 260, 220, "1", "System Prompt",
         "Agentens instruktionsbok: roll, ton, regler och vad agenten f\u00e5r eller inte f\u00e5r g\u00f6ra.")
    step(s, 340, 162, 260, 220, "2", "User Prompt",
         "Det anv\u00e4ndaren faktiskt skriver i chatten eller ber modellen om.")
    step(s, 620, 162, 240, 220, "3", "Kombineras",
         "B\u00e5da g\u00e5r in i samma anrop och blir en del av modellens arbetande kontext.")
    fnum(s, n)

    # ── 24. Kontextl\u00e4ngd ─────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Kontext", "Kontextl\u00e4ngd \u00e4r modellens korttidsminne", C.ACCENT2)
    rect(s, 60, 150, 820, 148, C.WHITE, C.BORDER)
    txt(s, 86, 172, 300, 24,
        "Allt detta t\u00e4vlar om samma budget", FT, 17, C.TEXT, bold=True)
    for x, w, label, color in [
        (90,  140, "Systemprompt",  C.ACCENT2),
        (240, 110, "Fr\u00e5ga",         C.ACCENT),
        (360, 126, "Historik",      C.MID_BLUE),
        (496, 156, "RAG / verktyg", C.LIGHT_BLU),
        (662, 158, "Svar",          C.DEEP_BLUE),
    ]:
        rect(s, x, 218, w, 30, color, rounded=False)
        txt(s, x, 224, w, 16, label, FB, 10.5, C.WHITE,
            bold=True, align=PP_ALIGN.CENTER)
    bullets(s, 60, 334, 760, 112, [
        "Kontextl\u00e4ngden \u00e4r den totala tokenbudgeten f\u00f6r hela anropet.",
        "B\u00e5de instruktioner, fr\u00e5ga, historik, RAG-data och svar f\u00e5r plats inom samma f\u00f6nster.",
        "Ju mer du stoppar in, desto mindre utrymme finns f\u00f6r annat.",
    ], 16)
    fnum(s, n)

    # ── 25. Resonerande modeller ──────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Kontext", "Resonerande modeller anv\u00e4nder ocks\u00e5 budget",
           C.ACCENT)
    txt(s, 60, 140, 340, 100,
        "I vissa resonerande modeller g\u00e5r en del av kontextbudgeten \u00e5t till mellanliggande resonemang, allts\u00e5 \"t\u00e4nkande\" som inte syns i det slutgiltiga svaret.",
        FB, 18, C.TEXT)
    rect(s, 430, 136, 470, 320, C.WHITE, C.BORDER)
    pic(s, "resoning", 448, 154, 434, 284)
    fnum(s, n)

    # ── 26. Context Rot ───────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Kontext", "Context Rot", C.ACCENT2)
    txt(s, 60, 140, 380, 180,
        "\u00c4ven om moderna modeller kan ha mycket l\u00e5nga kontextl\u00e4ngder betyder det inte att de anv\u00e4nder hela kontexten lika bra. N\u00e4r kontexten blir f\u00f6r l\u00e5ng, splittrad eller fylld av halvrelevant information f\u00e5r modellen sv\u00e5rare att hitta r\u00e4tt signal.",
        FB, 18, C.TEXT)
    rect(s, 470, 136, 430, 320, C.WHITE, C.BORDER)
    pic(s, "kontext", 488, 154, 394, 284)
    fnum(s, n)

    # ── 27. Kontexthantering ──────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Kontext", "Kontexthantering", C.ACCENT)
    step(s, 60, 160, 180, 220, "1", "Sammanfatta",
         "Kort ned \u00e4ldre historik n\u00e4r den inte l\u00e4ngre beh\u00f6vs i detalj.")
    step(s, 260, 160, 180, 220, "2", "Filtrera",
         "Skicka bara in det som faktiskt \u00e4r relevant f\u00f6r fr\u00e5gan.")
    step(s, 460, 160, 180, 220, "3", "Undvik brus",
         "Fyll inte kontexten med halvrelevant material bara f\u00f6r att det r\u00e5kar finnas.")
    step(s, 660, 160, 180, 220, "4", "Anv\u00e4nd RAG",
         "H\u00e4mta r\u00e4tt utdrag i st\u00e4llet f\u00f6r att skicka in allt.")
    fnum(s, n)

    # ── 28. RAG intro ─────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "RAG", "RAG ger modellen r\u00e4tt underlag i r\u00e4tt stund",
           C.ACCENT2)
    bullets(s, 60, 160, 320, 166, [
        "Minskar hallucinationer.",
        "Ger \u00e5tkomst till information som inte fanns i tr\u00e4ningsdatan.",
        "Skickar in relevanta utdrag i st\u00e4llet f\u00f6r hela dokument.",
    ], 18)
    rect(s, 420, 152, 200, 220, C.WHITE, C.BORDER)
    pic(s, "pdf", 438, 170, 164, 184)
    rect(s, 640, 152, 220, 220, C.WHITE, C.BORDER)
    pic(s, "nyheter", 658, 170, 184, 184)
    fnum(s, n)

    # ── 29. RAG: Chunks ───────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "RAG steg 1",
           "Ta in ett dokument och dela upp det i chunkar", C.ACCENT)
    rect(s, 60, 152, 240, 250, C.WHITE, C.BORDER)
    pic(s, "pdf", 82, 172, 196, 206)
    rect(s, 340, 152, 520, 250, C.WHITE, C.BORDER)
    pic(s, "chunks", 360, 170, 480, 214)
    fnum(s, n)

    # ── 30. RAG: Embeddings ───────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "RAG steg 2",
           "Embeddingsmodell: text blir s\u00f6kbar betydelse", C.ACCENT2)
    txt(s, 60, 140, 380, 80,
        "En embeddingsmodell l\u00e4ser varje stycke, f\u00f6rst\u00e5r betydelsen och skapar en vektorrepresentation som kan j\u00e4mf\u00f6ras semantiskt.",
        FB, 18, C.TEXT)
    rect(s, 60, 240, 840, 220, C.WHITE, C.BORDER)
    pic(s, "vektor", 80, 258, 800, 184)
    fnum(s, n)

    # ── 31. RAG: Retrieval ────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "RAG steg 3",
           "H\u00e4mta r\u00e4tt utdrag och skicka dem till modellen", C.ACCENT)
    step(s, 60, 152, 250, 120, "1", "Fr\u00e5gan \u2192 vektor",
         "\u00c4ven fr\u00e5gan g\u00f6rs om till en s\u00f6kbar betydelserepresentation.")
    step(s, 330, 152, 250, 120, "2", "N\u00e4rmaste hittas",
         "Systemet letar upp de mest relevanta chunkarna.")
    step(s, 600, 152, 250, 120, "3", "LLM svarar",
         "Utdragen skickas med fr\u00e5gan till spr\u00e5kmodellen.")
    rect(s, 60, 296, 790, 160, C.WHITE, C.BORDER)
    pic(s, "hitta", 80, 314, 750, 124)
    fnum(s, n)

    # ── 32. RAG: Begr\u00e4nsningar ───────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "RAG", "Begr\u00e4nsningar med RAG", C.ACCENT2)
    txt(s, 60, 140, 820, 60,
        "Spr\u00e5kmodellen ser bara de utdrag som retrieval-steget hittade \u2013 inte hela dokumentm\u00e4ngden. RAG fungerar bra f\u00f6r specifik information, men s\u00e4mre n\u00e4r fr\u00e5gan kr\u00e4ver:",
        FB, 20, C.TEXT)
    bullets(s, 60, 224, 820, 140, [
        "R\u00e4kning eller aggregering \u00f6ver m\u00e5nga delar av ett dokument.",
        "J\u00e4mf\u00f6relse mellan m\u00e5nga avsnitt samtidigt.",
        "F\u00f6rst\u00e5else av relationer utspridda \u00f6ver flera chunkar.",
    ], 18)
    rect(s, 60, 390, 820, 60, C.PANEL_ALT, C.BORDER, radius=0.12)
    txt(s, 82, 402, 776, 36,
        "Om svaret finns spritt \u00f6ver m\u00e5nga dokument finns risken att retrieval bara h\u00e4mtar en del \u2013 och modellen kan inte resonera \u00f6ver information den aldrig fick se.",
        FB, 13, C.MUTED)
    fnum(s, n)

    # ── 33. AI-agent intro ────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "AI-agent", "Vad \u00e4r en AI-agent?", C.ACCENT2)
    txt(s, 60, 148, 380, 128,
        "En AI-agent \u00e4r ett system d\u00e4r spr\u00e5kmodellen kan interagera med omv\u00e4rlden. Det handlar inte bara om att skriva text, utan om att kombinera modellen med verktyg, minne och flera steg \u00f6ver tid.",
        FB, 22, C.TEXT)
    rect(s, 470, 140, 430, 300, C.WHITE, C.BORDER)
    pic(s, "verktyg", 488, 158, 394, 264)
    fnum(s, n)

    # ── 34. Komponenter ───────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Komponenter", "Viktiga komponenter i en AI-agent", C.ACCENT)
    step(s, 60, 162, 250, 220, "1", "Verktyg",
         "L\u00e5ter modellen interagera med e-post, databaser, filer, kalendrar eller andra system.")
    step(s, 330, 162, 250, 220, "2", "Minne",
         "Sparar tidigare information s\u00e5 att agenten kan anv\u00e4nda den senare.")
    step(s, 600, 162, 250, 220, "3", "L\u00e4rande",
         "F\u00f6rb\u00e4ttring sker ofta genom minnen, feedback och orkestrering \u2013 snarare \u00e4n omtr\u00e4ning av sj\u00e4lva modellen.")
    fnum(s, n)

    # ── 35. Workflow vs Agent ─────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "J\u00e4mf\u00f6relse", "Workflow vs. Agent", C.ACCENT2)
    # Workflow
    rect(s, 60, 152, 400, 290, C.WHITE, C.BORDER)
    txt(s, 82, 168, 356, 22, "Workflow", FT, 16, C.ACCENT, bold=True)
    txt(s, 82, 196, 356, 36,
        "Styrt fl\u00f6de d\u00e4r processen redan \u00e4r f\u00f6rutbest\u00e4md.",
        FB, 14, C.TEXT)
    pic(s, "workflows", 82, 240, 356, 180)
    # Agent
    rect(s, 500, 152, 400, 290, C.WHITE, C.BORDER)
    txt(s, 522, 168, 356, 22, "Agent", FT, 16, C.ACCENT, bold=True)
    txt(s, 522, 196, 356, 36,
        "Modellen avg\u00f6r sj\u00e4lv vilka verktyg som ska anv\u00e4ndas och i vilken ordning.",
        FB, 14, C.TEXT)
    pic(s, "ai_agent", 522, 240, 356, 180)
    fnum(s, n)

    # ── 36. Agenttyper ────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    chrome(s, "Typer av agenter",
           "Konversationsagenter och autonoma agenter", C.ACCENT)
    step(s, 60, 162, 360, 240, "1", "Konversationsagent",
         "Pratar med en m\u00e4nniska och anv\u00e4nder verktyg p\u00e5 kommando. Bra n\u00e4r m\u00e4nniskan leder arbetet steg f\u00f6r steg.\n\nExempel: \"Boka ett m\u00f6te med Anna kl 14.\"")
    step(s, 460, 162, 360, 240, "2", "Autonom agent",
         "Kan triggas av h\u00e4ndelser eller scheman och arbeta i bakgrunden utan att en m\u00e4nniska beh\u00f6ver starta varje k\u00f6rning.\n\nExempel: En agent som \u00f6vervakar en inkorg dygnet runt.")
    fnum(s, n)

    # ── 37. Avslut ────────────────────────────────────────────────────
    n += 1
    s = new_slide(prs)
    rect(s, 0, 0, 960, 540, C.NAVY, rounded=False)
    rect(s, 0, 0, 960, 14, C.ACCENT2, rounded=False)
    txt(s, 60, 160, 840, 60, "Dags att b\u00f6rja bygga!", FT, 38, C.WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, 60, 240, 840, 60,
        "Nu n\u00e4r vi har begreppen p\u00e5 plats g\u00e5r vi vidare till n\u00e4sta kapitel och s\u00e4tter upp v\u00e5r milj\u00f6.",
        FB, 22, C.WHITE, align=PP_ALIGN.CENTER)
    txt(s, 60, 458, 840, 20, "Microsoft Copilot Studio",
        FB, 13, C.PANEL_ALT, bold=True, align=PP_ALIGN.CENTER)
    fnum(s, n)

    # ── Save ──────────────────────────────────────────────────────────
    prs.save(str(pptx_output))
    populate_media_slides(pptx_output)
    print(f"Created: {pptx_output}")
    print(f"Slides:  {n}")
    if make_pdf:
        export_pdf(pptx_output, pdf_output)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--pdf",
        action="store_true",
        help="Also export the generated presentation to PDF.",
    )
    parser.add_argument(
        "--media-pass",
        help="Internal helper mode: only embed video/audio into an existing presentation.",
    )
    args = parser.parse_args()
    if args.media_pass:
        raise SystemExit(0 if populate_media_slides(Path(args.media_pass), allow_external_python=False) else 1)
    build(make_pdf=args.pdf)
